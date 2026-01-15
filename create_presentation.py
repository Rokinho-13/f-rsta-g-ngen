"""Create a university-level PowerPoint presentation about the water cycle."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Vattnets Kretslopp"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Den Hydrologiska Cykeln"
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER

    # Academic info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Hydrologi & Geovetenskap | Universitetsnivå"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER


def add_overview_slide(prs):
    """Slide 2: Overview of the hydrological cycle"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Den Hydrologiska Cykeln - Översikt"
    p.font.size = Pt(32)
    p.font.bold = True

    # Content
    content = [
        "Definition: Kontinuerlig cirkulation av vatten genom atmosfären, hydrosfären, litosfären och biosfären",
        "",
        "Huvudprocesser:",
        "• Evaporation (E) - Övergång från vätska till gas",
        "• Transpiration (T) - Vattenavgivning från växter",
        "• Kondensation - Övergång från gas till vätska",
        "• Precipitation (P) - Nederbörd i alla former",
        "• Runoff (R) - Ytavrinning till vattendrag",
        "• Infiltration (I) - Nedträngning i marken",
        "",
        "Vattenbalanssekvation:",
        "P = E + T + R + ΔS    (där ΔS = förändring i lagring)"
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(6)


def add_evapotranspiration_slide(prs):
    """Slide 3: Evaporation and Transpiration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Evapotranspiration (ET)"
    p.font.size = Pt(32)
    p.font.bold = True

    content = [
        "Evaporation:",
        "• Energikrävande process: ~2,45 MJ/kg (latent värme)",
        "• Styrs av: strålning, temperatur, vindhastighet, luftfuktighet",
        "• Penman-ekvationen uppskattar potentiell evaporation",
        "",
        "Transpiration:",
        "• Biologisk process via växternas stomata",
        "• 95-99% av växternas vattenupptag transpireras",
        "• Regleras av: LAI (Leaf Area Index), markfuktighet, VPD",
        "",
        "Penman-Monteith ekvation (FAO-56):",
        "ET₀ = [0.408Δ(Rn-G) + γ(900/(T+273))u₂(es-ea)] / [Δ + γ(1+0.34u₂)]",
        "",
        "Globalt: ~70% av nederbörden återgår via ET"
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.space_after = Pt(5)


def add_precipitation_slide(prs):
    """Slide 4: Condensation and Precipitation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Kondensation & Precipitation"
    p.font.size = Pt(32)
    p.font.bold = True

    content = [
        "Kondensation:",
        "• Kräver kondensationskärnor (CCN) - aerosol, damm, salt",
        "• Sker när luften når daggpunkten (100% RH)",
        "• Adiabatisk kylning: ~6-10°C/1000m (torr/fuktig)",
        "",
        "Precipitationstyper:",
        "• Konvektiv - lokal uppvärmning, intensiv, kort duration",
        "• Orografisk - tvingad lyftning över berg",
        "• Frontal/Cyklonal - luftmassor möts vid fronter",
        "",
        "Mätning och analys:",
        "• Punktmätning: regnmätare (mm)",
        "• Spatial: radar (Z-R relation: Z = aRᵇ), satellit",
        "• IDF-kurvor: Intensity-Duration-Frequency för dimensionering",
        "",
        "Global medelårsnederbörd: ~1000 mm/år"
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.space_after = Pt(5)


def add_runoff_slide(prs):
    """Slide 5: Runoff and Groundwater"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Avrinning & Grundvatten"
    p.font.size = Pt(32)
    p.font.bold = True

    content = [
        "Ytavrinning (Runoff):",
        "• Hortonian overland flow - när P > infiltrationskapacitet",
        "• Saturation excess - när marken är vattenmättad",
        "• Avrinningskoefficient: C = R/P (varierar 0.1-0.9)",
        "",
        "Infiltration & Grundvatten:",
        "• Darcys lag: Q = -K·A·(dh/dl)",
        "• Hydraulisk konduktivitet (K): lera ~10⁻⁹ m/s, sand ~10⁻⁴ m/s",
        "• Akviferer: confined (artesisk) vs unconfined (fri yta)",
        "",
        "Uppehållstider i systemet:",
        "• Atmosfär: ~9 dagar",
        "• Floder: ~2 veckor",
        "• Grundvatten: 100-10 000 år",
        "• Oceaner: ~3 000 år",
        "",
        "Klimatförändringens påverkan: intensifierad cykel, extremer"
    ]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.space_after = Pt(4)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_overview_slide(prs)
    add_evapotranspiration_slide(prs)
    add_precipitation_slide(prs)
    add_runoff_slide(prs)

    output_path = "/home/user/f-rsta-g-ngen/vattnets_kretslopp.pptx"
    prs.save(output_path)
    print(f"Presentation sparad: {output_path}")


if __name__ == "__main__":
    main()
